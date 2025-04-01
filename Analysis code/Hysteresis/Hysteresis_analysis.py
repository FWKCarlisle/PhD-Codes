# -*- coding: utf-8 -*-
"""
Created on Mon Feb 17 10:49:54 2025
February 2025 beamtime XMCD analysis codes
@author: Wstev
"""
import matplotlib.pyplot as plt
import numpy as np
from nexusformat.nexus import nxload
from pathlib import Path
from scipy.integrate import simpson 
from mpl_toolkits.axes_grid1.inset_locator import inset_axes

import datetime
import csv
import pptx
import pptx.util
import os

class Nexus_handling:
    def get_nexus_data_I06(self,file):
        """Function that loads the data from a nexus file and returns it as a list of numpy
arrays"""
        entry_string = "entry"
        data_region_list = []
        metadata_region_list = None
        x_array = file[entry_string]["instrument"]["fastEnergy"]["value"].nxvalue
        y_array = file[entry_string]["instrument"]["fesData"]["C1"].nxvalue
        data_region_list.append({"x": x_array, "y": y_array})
        y_array = file[entry_string]["instrument"]["fesData"]["C5"].nxvalue
        data_region_list.append({"x": x_array, "y": y_array})
        y_array = file[entry_string]["instrument"]["fesData"]["idio"].nxvalue
        data_region_list.append({"x": x_array, "y": y_array})
        y_array = file[entry_string]["instrument"]["fesData"]["ifiofb"].nxvalue
        data_region_list.append({"x": x_array, "y": y_array})
        magnet_field = file[entry_string]["instrument"]["scm"]["field_z"].nxvalue
        metadata_region_list = {"magnet_ﬁeld": magnet_field}
        polar = file[entry_string]["instrument"]["id"]["polarisation"].nxvalue
        metadata_region_list["polarisation"] = polar
        temp = file[entry_string]["instrument"]["scm"]["T_sample"].nxvalue
        metadata_region_list["T_sample"] = temp
        ###Add temperature values here
        ### Add angle of the beam here => 
        return data_region_list, metadata_region_list

    def open_single_spectra(self,ﬁle_number,directory_path,ﬁle_preﬁx,sensor):
        ﬁle_name = directory_path + ﬁle_preﬁx + str(ﬁle_number) + ".nxs"
        spectra_ﬁle = nxload(ﬁle_name)
        data,meta_data = self.get_nexus_data_I06(spectra_ﬁle)
        if sensor == "TEY":
            spectra = {"x" : data[2]["x"],"y":data[2]["y"],"meta" : meta_data}
        elif sensor == "TFY":
            spectra = {"x" : data[3]["x"],"y":data[3]["y"],"meta" : meta_data}
        return spectra

class XMCD_data_analysis(Nexus_handling):
    def __init__(self):
        self.on_oI_edge_points = [571, 576.9]
        self.directory_path = None
        self.ﬁle_preﬁx = None
        self.energy = None
        self.spectra = []
        self.hyst = []
        self.XMCD = []
        self.magz = []
        #self.temperature
   

    def set_directory_path(self, directory_path):
        self.directory_path = directory_path

    def set_ﬁle_preﬁx(self, ﬁle_preﬁx):
        self.ﬁle_preﬁx = ﬁle_preﬁx

    def load_spectra(self, directory_path, ﬁle_preﬁx, XMCD_spectra_Set,sensor = "TEY"):  
        self.directory_path = directory_path
        self.ﬁle_preﬁx = ﬁle_preﬁx

        if len(XMCD_spectra_Set) != 1:
            for ﬁle_number in XMCD_spectra_Set: 
                print(file_number)
                data_set = self.open_single_spectra(ﬁle_number, directory_path,  ﬁle_preﬁx,sensor)
                self.spectra.append(data_set)
                self.energy = self.spectra[0]["x"]
        else:
            print(XMCD_spectra_Set)
            data_set = self.open_single_spectra(XMCD_spectra_Set[0], directory_path,  ﬁle_preﬁx,sensor)
            self.spectra.append(data_set)
    
            self.energy = self.spectra[0]["x"]
#######################################################################################
    ### Plotting functions
    
    def plot_spectra(self,XAS_file_Set,i, bounds = None, normalisation = None):
        #Plots a single XAS spectra out of a set 
        #XAS_file_set : a set of XAS scan numbers
        #i : the indice of which scan you want to plot
        Energy,Idio = self.XMCD.spectra[i]["x"],self.XMCD.spectra[i]["y"]
        if bounds is not None:
            plt.xlim(bounds[0])
            plt.ylim(bounds[1])

        if normalisation is not None:
            value = np.mean(Idio[:normalisation])
            plotting_y = Idio - value
            plotting_y = plotting_y / np.max(plotting_y) 
        else:
            plotting_y = Idio


        plt.plot(Energy,plotting_y,label = XAS_file_Set[i])
        plt.title(("XAS ",sensor,XAS_file_set[i],self.spectra[i]["meta"]["polarisation"]))
        plt.legend()
        plt.xlabel("Energy(eV)")
        plt.ylabel("idio")   

    def Dichroism(self,XMCD_ﬁle_Set, normalisation = None,scaling_factor=1, avg = "Before", on_top = False):
        #Takes in a series of XAS spectra of opposing polarisations then averages all the nc and all the pc and subtracts them to give us the dichroism
        # XMCD_file_Set : list of XAS spectra, at a minimum we need one of each polarisation of light.

        fig = plt.figure(figsize=(12, 6))
        
        if on_top:
            [axXAS, axXMCD] = fig.subplots(1,2, sharey=True)  # (nrows, ncols, index)
            
        else:
            [axXAS, axXAS_pc, axXMCD] = fig.subplots(1,3, sharey=True)  # (nrows, ncols, index)
            axXAS_pc.set_title("pc")
        
        
        axXAS.set_title("nc")
        axXMCD.set_title("XMCD")
        
        nc_storage = np.zeros(len(self.energy))
        pc_storage = np.zeros(len(self.energy))
        
        B_Field = self.spectra[0]["meta"]["magnet_ﬁeld"]
        temp = self.spectra[0]["meta"]["T_sample"]
        fig.suptitle("Dichroism at " + str(np.round(B_Field,2)) +  f"T {round(temp,3)}K")
        print("B Field - ",B_Field)
        
        
        
        axXMCD.set_xlabel("energy (eV")
        axXAS.set_ylabel("idio")
        
        n,m = 0,0
        max_point = 0

        
        for i in range(0,len(XMCD_file_Set)):
            Polarisation = self.spectra[i]["meta"]["polarisation"]
            F = self.spectra[i]["y"]

            if normalisation is not None:
                value = np.mean(F[:normalisation])
                F -= value

            
            f=np.polyfit(self.energy, F, 1)

            fit=np.polyval(f,self.energy)
            mx = np.max((F-fit)+0.1)
            if mx > max_point:
                max_point = mx
            if Polarisation == "nc":
                F = F * scaling_factor
                nc_storage += F
                n += 1
                axXAS.plot(self.energy,F, label = f"{XMCD_file_Set[i]} nc" )
                axXAS.legend()

            elif Polarisation == "pc":
                pc_storage = pc_storage + F
                m = m + 1
                if on_top:
                    axXAS.plot(self.energy,F,label = f"{XMCD_file_Set[i]}  pc" )
                else:
                    axXAS_pc.plot(self.energy,F,label = XMCD_file_Set[i] )
                    axXAS_pc.legend()

        if avg == "Before":
            nc_storage = nc_storage/n
            pc_storage = pc_storage/m
            dichroism = abs(pc_storage) - abs(nc_storage)

            max_dichroism = np.max(dichroism)
            percentage = 100*max_dichroism/max_point
            print(f"Max Di - {max_dichroism}, Max Signal - {max_point}")
        
        axXMCD.plot(self.energy,dichroism,label = f"XMCD - {np.round(percentage,3)}%")
        axXMCD.legend()

        
        axXMCD.axhline(y = 0, color = 'b', linestyle = 'dashed') 
        axXAS.axhline(y = 0, color = 'b', linestyle = 'dashed') 
        if not on_top:
            axXAS_pc.axhline(y = 0, color = 'b', linestyle = 'dashed') 
            
        return self.energy,dichroism                    

def prelimenary_code(directory_path, ﬁle_preﬁx, XAS_ﬁle_set,sensor):
    XMCD = XMCD_data_analysis()
    XMCD.load_spectra(directory_path, ﬁle_preﬁx, XAS_ﬁle_set,sensor)
    return XMCD

def get_values(file):
    
    nexus = nxload(file)
    TEY_A=nexus.entry.instrument.hyst2.detector1_A.nxvalue 
    TEY_B=nexus.entry.instrument.hyst2.detector1_B.nxvalue 
    I_0_A = nexus.entry.instrument.hyst2.detector2_A.nxvalue 
    I_0_B = nexus.entry.instrument.hyst2.detector2_B.nxvalue 
    
    return [TEY_A, TEY_B, I_0_A, I_0_B]



def Hysteresis(files, offset=[0,0]):
    """

    Parameters
    ----------
    files : Array
        List of nexus files generated from fast hysteris code.
    offset : Array, optional
        DESCRIPTION. The default is [0,0]. This is an [Y,X] offset for each arm of the loop

    Returns
    -------
        This plots the hysterisis loops that used only a single on and off resonance point.It also allows the loops to
        be offset from each other. It plots both the raw data and the data normalised with the IO.
        
    """
    if len(files) != 4:
        print("Wrong number of files")
        return 0
    values = {}
    
    for file in files:
        data = get_values("FILES/HOPG1/Hysteresis_analysis/Hysteresis/" + file)
        for i, datum in enumerate(data):
            values[f"{file} {i}"] = datum # 0 = TeyA, 1 = TEYB, 2 = I0A, 3 = I0B

    nexus = nxload("FILES/HOPG1/Hysteresis_analysis/Hysteresis/" + files[0])
    B_field = nexus.entry.instrument.hyst2.value.nxvalue
    temp = nexus.entry.instrument.scm.T_sample
    
    
    norm_dict = {}
    numerator_6_8 = (((values[f"{files[0]} 0"] / values[f"{files[0]} 2"]) - (values[f"{files[0]} 1"] / values[f"{files[0]} 3"])) - ((values[f"{files[2]} 0"] / values[f"{files[2]} 2"]) - (values[f"{files[2]} 1"] / values[f"{files[2]} 3"])))
    denominator_6_8 = 0.5* (((values[f"{files[0]} 0"] / values[f"{files[0]} 2"]) - (values[f"{files[0]} 1"] / values[f"{files[0]} 3"])) + ((values[f"{files[2]} 0"] / values[f"{files[2]} 2"]) - (values[f"{files[2]} 1"] / values[f"{files[2]} 3"])))
    plotting_1 = numerator_6_8/denominator_6_8

    numerator_7_9 = -(((values[f"{files[1]} 0"] / values[f"{files[1]} 2"]) - (values[f"{files[1]} 1"] / values[f"{files[1]} 3"])) - ((values[f"{files[3]} 0"] / values[f"{files[3]} 2"]) - (values[f"{files[3]} 1"] / values[f"{files[3]} 3"])))
    denominator_7_9 = 0.5* (((values[f"{files[1]} 0"] / values[f"{files[1]} 2"]) - (values[f"{files[1]} 1"] / values[f"{files[1]} 3"])) + ((values[f"{files[3]} 0"] / values[f"{files[3]} 2"]) - (values[f"{files[3]} 1"] / values[f"{files[3]} 3"])))
    plotting_2 = numerator_7_9/denominator_7_9

    plt.figure(1)
    plt.plot(B_field, numerator_6_8, "s-" ,label = "PC") # just (A - B) - (C-D)
    plt.plot(B_field, numerator_7_9,"d-" ,label = "nC")
    plt.title(("Raw data",temp,"TEY"))
    plt.axhline(y=0)
    plt.axvline(x=0)
    plt.legend()
    
    plt.figure(2)
    plt.plot(B_field, plotting_1+offset[0],"s-", label = "PC_norm") # (A - B) - (C-D) / (0.5*(A - B) + (C-D))
    plt.plot(B_field, plotting_2+offset[1],"d-", label = "nC_norm")#
    plt.title(("Normalised data",temp,"TEY"))
    plt.legend()
    plt.axhline(y=0)
    plt.axvline(x=0)

def find_nearest(array, value):
    array = np.asarray(array)
    idx = (np.abs(array - value)).argmin()
    return array[idx], idx

def method1(XMCD):
    XMCD,pc_Bfield, nc_Bfield,pc_Fenergy, temp, label = XMCD

    hyst_point = max(XMCD)-min(XMCD)
    Bfield = (pc_Bfield + nc_Bfield)/2
    
    return hyst_point, Bfield, temp

def method2(XMCD):
    XMCD,pc_Bfield, nc_Bfield,pc_Fenergy, temp, label = XMCD

    avg_bkgd = np.mean(XMCD[:10])
    nearest_point, idx = find_nearest(pc_Fenergy, 576.8)
    hyst_point = XMCD[idx] - avg_bkgd
    Bfield = (pc_Bfield + nc_Bfield)/2
    
    return hyst_point, Bfield, temp

def method3(XMCD):
    XMCD,pc_Bfield, nc_Bfield,pc_Fenergy, temp, label = XMCD

    xxxx, idx_max = find_nearest(pc_Fenergy, 576.8)
    xxxx, idx_min = find_nearest(pc_Fenergy, 578)
    max_point = XMCD[idx_max]
    min_point = XMCD[idx_min]

    hyst_point = max_point - min_point
    Bfield = (pc_Bfield + nc_Bfield)/2
    
    return hyst_point, Bfield, temp

def get_remenance(hyst, B):

    idx_list = []
    
    for i, x in enumerate(B):
        print(x)
        if abs(x) == 0.0:
            idx_list.append(i)
    
    i_1 = idx_list[0]
    i_2 = idx_list[1]
    
    rem_1 = hyst[i_1]
    rem_2 = hyst[i_2]


    print(f"Rem 1 ({B[i_1]} ): {rem_1}, Rem 2 ({B[i_2]} ): {rem_2}, half_point: {0.5*(abs(rem_1) + abs(rem_2))}, full: {0.5*(abs(rem_1) + abs(rem_2)) - min(rem_1, rem_2)}")

    
    remenance = 0.5*(abs(rem_1) + abs(rem_2)) 

    return remenance
    
    

def manual_hyst(data_dir,Yield = "TEY",file_name = "Manual_Hysteresis_1", normalisation=None,bounds =None, plot_xmcd=True, Label = False, scaling = None):
    """
    

    Parameters
    ----------
    data_dir : STRING
        File path where the spectra are located.
    Yield : String, optional
        DESCRIPTION. The default is "TEY". Decides whether to use total elecron yield (TEY) or total fluorescencse yield (TFY)
    normalisation : TYPE, optional
        DESCRIPTION. Unsure what this does
    bounds : TYPE, array
        DESCRIPTION. The default is None, Sets X,Y bounds for plotting
    plot_xmcd : TYPE, optional
        DESCRIPTION. The default is True. Chooses whether to plot each spectra that contributes a point to the hysterisis loop
    Label : TYPE, optional
        DESCRIPTION. The default is False.
    scaling : TYPE, optional
        DESCRIPTION. If true this value scales the nc polarisations by an arbitrary factor.

    Returns
    -------
    None.

    """
    files = [f.name for f in data_dir.iterdir() if f.is_file()]
    files.sort()

    mid = int(len(files)/2)
    
    pc_files = files[:mid]
    nc_files = files[mid:]
    # print("PC: ", pc_files)
    # print("NC: ", nc_files)

    final_pc = []
    final_nc = []
    
    for i, file in enumerate(files):
        nexus = nxload(str(data_dir)+ "/" + file)
        B_field = nexus.entry.instrument.scm.field_z.nxvalue
        if Yield == "TEY":
            TEY=nexus.entry.instrument.fesData.C1.nxvalue 
        elif Yield == "TFY":
            TEY=nexus.entry.instrument.fesData.C5.nxvalue     
        I0=nexus.entry.instrument.fesData.C2.nxvalue
        Fenergy = nexus.entry.instrument.fastEnergy.value.nxvalue
        polar = nexus.entry.instrument.id.polarisation.nxvalue
        temp = nexus.entry.instrument.scm.T_sample

        norm_TEY = TEY/I0

        if normalisation is not None:
            value = np.mean(norm_TEY[:normalisation])
            TEY_final = norm_TEY - value
        else:
            TEY_final = norm_TEY

        data = [TEY_final, B_field, Fenergy,temp, file]
        
        if polar == "pc":
            final_pc.append(data)
            # print(data[3], i)
        elif polar == "nc":
            final_nc.append(data)
            # print("nc",data[3], i)

    
    XMCDs = []
    integrations = []
    if os.path.exists(rf"PhD-Codes\Analysis code\Hysteresis\{file_name}.pptx"):
        os.remove(rf"PhD-Codes\Analysis code\Hysteresis\{file_name}.pptx")
    prs = setup_presentation()
    for i, data in enumerate(final_pc):
        pc_TEY, pc_Bfield, pc_Fenergy, pc_temp, pc_file = data
        nc_TEY, nc_Bfield, nc_Fenergy, nc_temp, nc_file = final_nc[i]
        if scaling is not None:
            nc_TEY *= scaling
        XMCD = pc_TEY - nc_TEY
        height_pc = max(pc_TEY) - min(pc_TEY)
        height_nc = max(nc_TEY) - min(nc_TEY)
        XMCD_norm = XMCD / (0.5 * (height_pc + height_nc))

        XMCD_change = np.mean(XMCD_norm[:10])
        XMCD_norm -= XMCD_change

        # Square the XMCD_norm values
        XMCD_norm_squared = XMCD_norm**2

        # Perform numerical integration using simps
        integration_result = np.sqrt(simpson(y=XMCD_norm_squared, x=pc_Fenergy))
        integrations.append(integration_result)
        label = f"{pc_file} - {nc_file} @ {pc_Bfield}/{nc_Bfield} @ {temp}K"
        # print(label)
        temp =(pc_temp + nc_temp)*0.5
        
        XMCD_data = [XMCD_norm, pc_Bfield, nc_Bfield,  pc_Fenergy,temp, label]

        percentage = 100*(max(abs(XMCD))/max((max(abs(nc_TEY)) ,max(abs(pc_TEY)))))
        
        XMCDs.append(XMCD_data)
        if plot_xmcd:
            plt.figure(i)
            
            fig, [axXAS, axXMCD] = plt.subplots(1,2, sharey=True)
            axXAS.plot(pc_Fenergy, pc_TEY, label="pc")
            axXAS.plot(nc_Fenergy, nc_TEY, label="nc")
            axXAS.set_title("XAS")
            axXAS.legend()

            # Plot XMCD data
            axXMCD.plot(pc_Fenergy, XMCD_norm, label="XMCD_norm")
            axXMCD.plot(pc_Fenergy, XMCD_norm_squared, label="XMCD_norm_sqrd")
            axXMCD.fill_between(pc_Fenergy, 0, XMCD_norm_squared, color='orange', alpha=0.3, label="Integration Area")
            axXMCD.set_title("XMCD")
            axXMCD.legend()

            # Add integration result to the plot
            axXMCD.text(0.05, 0.95, f"Integration: {integration_result:.2f}", transform=axXMCD.transAxes,
                fontsize=10, verticalalignment='top', bbox=dict(boxstyle="round", facecolor="white", alpha=0.5))

            plt.title(label)
            
            plt.axhline(y = 0, color = 'b', linestyle = 'dashed', label=f"{ round(percentage,3) }% ")
            plt.legend()
            # plt.show()

            plt.savefig(r"PhD-Codes\Analysis code\Hysteresis\data_for_ppt\plot.png")

            save_to_presentation(prs, r"PhD-Codes\Analysis code\Hysteresis\data_for_ppt\plot.png",filename = file_name)
    # print(len(XMCDs))
    plt.close("all")
    plt.figure(len(final_pc)+2) 
    M1_hyst, M1_B = [], []
    M2_hyst, M2_B= [], []
    M3_hyst, M3_B = [], []
    temps = []
    for j, XMCD in enumerate(XMCDs):
        # hyst, B, temp = method1(XMCD)
        # M1_hyst.append(hyst)
        # M1_B.append(B)

        # hyst, B, temp = method2(XMCD)
        # M2_hyst.append(hyst)
        # M2_B.append(B)
        
        hyst, B, temp = method3(XMCD)
        M3_hyst.append(hyst)
        M3_B.append(B)
        temps.append(temp)
        if Label:
            plt.text(B,hyst, f"{j} {60-j}")
        
    co = np.arange(0,len(XMCDs),1)

    for i in range(len(M3_B)):
        if M3_B[i] < 0:
            integrations[i] = -integrations[i]

    
    plt.figure(i+1)
    
    plt.scatter(M3_B, M3_hyst,marker='o' ,c=co,cmap='inferno', label = "Raw XMCD")
    plt.scatter(M3_B, integrations,marker='x',c=co,cmap='inferno', label = "Integration")
    
    if bounds is not None:
        plt.xlim(bounds[0])
        plt.ylim(bounds[1])
    
    plt.xlabel("$B_z (T)$")
    plt.ylabel("$XMCD_{Norm}$")

    data_dir = str(data_dir)
    plt.axhline(y=0)
    plt.axvline(x=0)
    plt.title((data_dir[len(data_dir)-20:len(data_dir)],"$T_{average} = $",round(np.mean(temps),5),"K"))
   
    ax_inset = inset_axes(plt.gca(), width="40%", height="40%", loc="lower right")  # Adjust size and location
    ax_inset.scatter(M3_B, M3_hyst, marker='o', c=co, cmap='inferno')
    ax_inset.scatter(M3_B, integrations, marker='x', c=co, cmap='inferno')

    inset_bounds = ([-0.04, 0.04], [min(M3_hyst)-0.1,max(M3_hyst)+0.1])  # Example bounds for the inset
    ax_inset.set_xlim(inset_bounds[0])
    # ax_inset.set_ylim(inset_bounds[1])

    ax_inset.set_title("Zoomed Inset", fontsize=8)
    ax_inset.tick_params(axis='both', which='major', labelsize=6)

    plt.colorbar(label = "Index")
    plt.savefig(r"PhD-Codes\Analysis code\Hysteresis\data_for_ppt\plot.png")

    save_to_presentation(prs, r"PhD-Codes\Analysis code\Hysteresis\data_for_ppt\plot.png",filename = file_name)
    # plt.show()
    # remenance =  get_remenance(M3_hyst, M3_B)
    plt.close("all")
    return max(M3_hyst)#, remenance, remenance/max(M3_hyst)

def setup_presentation():
    prs = pptx.Presentation()
    return prs

def save_to_presentation(prs,image_path, filename = "presentation"):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Step 4: Add image to the slide
    left = pptx.util.Inches(0.25)   # Adjust positioning
    top = pptx.util.Inches(1.5)
    width = pptx.util.Inches(8)  # Set width (optional)
    height = pptx.util.Inches(6) # Set height (optional)
    slide.shapes.add_picture(image_path, left, top, width, height)

    # Step 5: Save the PowerPoint file
    pptx_filename = rf"PhD-Codes\Analysis code\Hysteresis\{filename}.pptx"
    prs.save(pptx_filename)

    # print(f"Presentation saved as {pptx_filename}")


print("Saved Code")
directory_path = r"C:\Users\ppxfc1\OneDrive - The University of Nottingham\Desktop\PhD\CrCl3\MM38256-1\HOPG1\Hysteresis_analysis"
ﬁle_preﬁx = "i06-1-"
in_bounds = [[-0.1,0.1],[-0.1,0.1]]
print("ALL GOOD HERE BOSS: ", datetime.datetime.now())
# list = ["2k","2k_2","2k_GI","2k_NI","8k","10k","11.8k","11k","12.2k","12.4","12.5k","12.5k_2","12.6_2",]
# list = ["BP_12.4k","12.6k","12.7k","12.8k","12.8k_2","12.9k","12k","12k_2","13.2k","13k","13k_2","14k","15k","17k", ]

list = ["1","2","3","4","5","6","7","10k","11k","12.5k","12k","12k_2","13.5k","13k","14k","14k_2","20k",]
for i, x in enumerate(list):
    if x[0] == "B":
        file_name = x
    else:
        # file_name = f"Hyst_{x}"
        file_name = f"Manual_Hysteresis_{x}"
    file_path = Path(directory_path + f"/{file_name}/")
    manual_hyst(file_path,plot_xmcd=True,bounds = ([-2.1,2.1],[-1.5,1.5]), file_name=file_name)





