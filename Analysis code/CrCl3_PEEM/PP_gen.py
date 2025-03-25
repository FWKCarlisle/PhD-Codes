import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Step 1: Generate a Matplotlib figure
fig, ax = plt.subplots(2, 2, figsize=(8, 6))  # Example 2x2 subplots
for i in range(2):
    for j in range(2):
        ax[i, j].plot([0, 1, 2], [0, i+j, (i+j)**2])  # Simple plots
plt.tight_layout()

# Step 2: Save the figure as an image
image_path = "plot.png"
fig.savefig(image_path, dpi=300)  # Save as high-resolution PNG
plt.close(fig)  # Close the figure to free memory

# Step 3: Create a PowerPoint Presentation and add the image
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # Use a blank slide
slide = prs.slides.add_slide(slide_layout)

# Step 4: Add image to the slide
left = Inches(1)   # Adjust positioning
top = Inches(1)
width = Inches(8)  # Set width (optional)
height = Inches(6) # Set height (optional)
slide.shapes.add_picture(image_path, left, top, width, height)

# Step 5: Save the PowerPoint file
pptx_filename = "PhD-Codes\Analysis code\CrCl3_PEEM\presentation.pptx"
prs.save(pptx_filename)

print(f"Presentation saved as {pptx_filename}")
